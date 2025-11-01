"""PowerPoint 解析器"""

import io
import logging

from pptx import Presentation

from app.core.parsers.base import BaseParser

logger = logging.getLogger(__name__)


class PowerPointParser(BaseParser):
    """PowerPoint 文档解析器"""

    async def parse(self, file_data: bytes, **_kwargs) -> str:
        """解析 PowerPoint 文档"""
        try:
            prs = Presentation(io.BytesIO(file_data))

            text_parts = []

            for i, slide in enumerate(prs.slides, 1):
                text_parts.append(f"Slide {i}:\n")

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text)

                text_parts.append("\n")

            return "\n".join(text_parts)

        except Exception as e:
            logger.error(f"Failed to parse PowerPoint: {e}")
            raise
